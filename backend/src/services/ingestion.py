import hashlib
import re
import json
import asyncio
from io import BytesIO
from xml.etree import ElementTree
from zipfile import BadZipFile, ZipFile
from logging import getLogger

from pypdf import PdfReader
from PIL import Image
import pytesseract
from pptx import Presentation
import openpyxl
from bs4 import BeautifulSoup

from src.services.ollama import OllamaClient

logger = getLogger(__name__)

_SENTENCE_SPLIT = re.compile(r"(?<=[.!?])\s+(?=[A-Z0-9\"'\(])")
_PARAGRAPH_SPLIT = re.compile(r"\n\s*\n+")

SUPPORTED_TEXT_MIME_PREFIXES = ("text/",)
SUPPORTED_TEXT_MIME_TYPES = {
    "application/json",
    "application/markdown",
    "application/x-ndjson",
    "text/markdown",
    "text/csv",
}
DOCX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def parse_docx_bytes(content: bytes) -> str:
    """Extract plain text from a DOCX file's XML structure."""
    try:
        with ZipFile(BytesIO(content)) as archive:
            document_xml = archive.read("word/document.xml")
    except (BadZipFile, KeyError) as exc:
        raise ValueError("Could not read DOCX document text") from exc

    root = ElementTree.fromstring(document_xml)
    paragraphs: list[str] = []
    namespace = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    for paragraph in root.findall(".//w:p", namespace):
        text = "".join(node.text or "" for node in paragraph.findall(".//w:t", namespace)).strip()
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def sha256_bytes(content: bytes) -> str:
    """Generate a SHA256 hash of raw bytes for document deduplication."""
    return hashlib.sha256(content).hexdigest()


def parse_pdf(content: bytes) -> str:
    reader = PdfReader(BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    text = "\n\n".join(page.strip() for page in pages if page.strip())
    return text


def parse_image_ocr(content: bytes) -> str:
    image = Image.open(BytesIO(content))
    return pytesseract.image_to_string(image)


def parse_pptx(content: bytes) -> str:
    prs = Presentation(BytesIO(content))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


def parse_xlsx(content: bytes) -> str:
    wb = openpyxl.load_workbook(BytesIO(content), data_only=True)
    text = []
    for sheet in wb.worksheets:
        for row in sheet.iter_rows(values_only=True):
            row_vals = [str(cell) for cell in row if cell is not None]
            if row_vals:
                text.append(" ".join(row_vals))
    return "\n".join(text)


def parse_html(content: bytes) -> str:
    soup = BeautifulSoup(content, "html.parser")
    return soup.get_text(separator="\n", strip=True)


def parse_text(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


def parse_document_bytes(content: bytes, mime_type: str | None, filename: str) -> str:
    normalized_mime = (mime_type or "").lower()
    suffix = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    text = ""
    try:
        if normalized_mime == "application/pdf" or suffix == "pdf":
            text = parse_pdf(content)
        elif normalized_mime == DOCX_MIME_TYPE or suffix == "docx":
            text = parse_docx_bytes(content)
        elif suffix == "pptx" or "presentation" in normalized_mime:
            text = parse_pptx(content)
        elif suffix == "xlsx" or "spreadsheet" in normalized_mime:
            text = parse_xlsx(content)
        elif suffix in {"html", "htm"} or "html" in normalized_mime:
            text = parse_html(content)
        elif suffix in {"jpg", "jpeg", "png", "tiff"} or normalized_mime.startswith("image/"):
            text = parse_image_ocr(content)
        elif normalized_mime.startswith(SUPPORTED_TEXT_MIME_PREFIXES) or normalized_mime in SUPPORTED_TEXT_MIME_TYPES:
            text = parse_text(content)
        elif suffix in {"txt", "md", "markdown", "csv", "json"}:
            text = parse_text(content)
        else:
            raise ValueError(f"Unsupported document type: {mime_type or suffix or 'unknown'}")
    except Exception as e:
        raise ValueError(f"Error parsing document: {e}")

    text = "\n".join(line.rstrip() for line in text.splitlines()).strip()
    if not text:
        raise ValueError("Document did not contain extractable text")
    return text


def chunk_text(text: str, chunk_words: int = 260, overlap_words: int = 40) -> list[str]:
    """Simple sliding-window split of text into word-based chunks."""
    words = text.split()
    if not words:
        return []

    chunks: list[str] = []
    step = max(chunk_words - overlap_words, 1)
    for start in range(0, len(words), step):
        chunk = " ".join(words[start : start + chunk_words]).strip()
        if chunk:
            chunks.append(chunk)
        if start + chunk_words >= len(words):
            break
    return chunks


def _split_sentences(text: str) -> list[str]:
    """Split text into individual sentences using a regex boundary."""
    sentences: list[str] = []
    for sent in _SENTENCE_SPLIT.split(text):
        s = sent.strip()
        if s:
            sentences.append(s)
    return sentences


def _pack_units(units: list[str], target_words: int) -> list[str]:
    """Greedily pack ordered text units into chunks up to target_words."""
    packed: list[str] = []
    buffer: list[str] = []
    buffer_words = 0
    for unit in units:
        unit_words = len(unit.split())
        if buffer and buffer_words + unit_words > target_words:
            packed.append(" ".join(buffer))
            buffer = []
            buffer_words = 0
        buffer.append(unit)
        buffer_words += unit_words
    if buffer:
        packed.append(" ".join(buffer))
    return packed


def chunk_markdown(text: str, child_words: int = 80) -> list[tuple[str, str]]:
    """Splits markdown content based on headings, prepending the heading hierarchy to chunks."""
    lines = text.split("\n")
    sections = []
    current_headings = []
    current_section_lines = []
    
    for line in lines:
        stripped = line.strip()
        header_match = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if header_match:
            if current_section_lines:
                heading_path = " > ".join(h[1] for h in current_headings)
                sections.append({
                    "heading_path": heading_path,
                    "content": "\n".join(current_section_lines)
                })
                current_section_lines = []
            
            level = len(header_match.group(1))
            title = header_match.group(2).strip()
            
            while current_headings and current_headings[-1][0] >= level:
                current_headings.pop()
            current_headings.append((level, title))
        else:
            current_section_lines.append(line)
            
    if current_section_lines:
        heading_path = " > ".join(h[1] for h in current_headings)
        sections.append({
            "heading_path": heading_path,
            "content": "\n".join(current_section_lines)
        })
        
    result = []
    for sec in sections:
        content_text = sec["content"].strip()
        if not content_text:
            continue
        sentences = _split_sentences(content_text)
        children = _pack_units(sentences, child_words)
        
        prefix = f"[Context: {sec['heading_path']}]\n" if sec["heading_path"] else ""
        for child in children:
            if child.strip():
                child_with_context = prefix + child
                parent_content = f"# {sec['heading_path']}\n{content_text}" if sec["heading_path"] else content_text
                result.append((child_with_context, parent_content))
                
    return result


def chunk_text_hierarchical(
    text: str,
    child_words: int = 80,
    parent_words: int = 320,
    filename: str = "",
) -> list[tuple[str, str]]:
    """Implement structure-aware semantic hierarchical chunking by file type."""
    text = text.strip()
    if not text:
        return []

    suffix = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

    # 1. Spreadsheet Chunking (Excel / CSV)
    if suffix in {"xlsx", "csv"}:
        lines = [line.strip() for line in text.split("\n") if line.strip()]
        if len(lines) <= 1:
            return [(text, text)]
        
        header = lines[0]
        data_rows = lines[1:]
        
        # Group rows together (15 rows per child chunk)
        chunk_size = 15
        result = []
        for i in range(0, len(data_rows), chunk_size):
            row_group = data_rows[i : i + chunk_size]
            child_text = f"Column Headers: {header}\nRows:\n" + "\n".join(row_group)
            
            # Parent context includes adjacent row ranges
            parent_group = data_rows[max(0, i - chunk_size) : min(len(data_rows), i + 2 * chunk_size)]
            parent_text = f"Column Headers: {header}\nRows context:\n" + "\n".join(parent_group)
            
            result.append((child_text, parent_text))
        return result

    # 2. Markdown Chunking
    if suffix in {"md", "markdown"}:
        return chunk_markdown(text, child_words)

    # 3. Default Hierarchical Chunking (PDF, Docx, OCR, plain text)
    paragraphs = [p.strip() for p in _PARAGRAPH_SPLIT.split(text) if p.strip()]
    if not paragraphs:
        paragraphs = [text]

    parents = _pack_units(paragraphs, parent_words)

    result: list[tuple[str, str]] = []
    seen_children: set[str] = set()
    for parent_text in parents:
        sentences = _split_sentences(parent_text)
        if not sentences:
            continue
        children = _pack_units(sentences, child_words)
        for child in children:
            if child and child not in seen_children:
                seen_children.add(child)
                result.append((child, parent_text))

    return result


async def embed_chunks(
    chunks: list[tuple[str, str]],
    embed_model: str | None = None,
) -> list[tuple[str, list[float], dict]]:
    """Legacy helper: embeds child text directly without question enrichment."""
    ollama = OllamaClient()
    embedded: list[tuple[str, list[float], dict]] = []
    for index, (child, parent) in enumerate(chunks):
        embedding = await ollama.embed(child, model_name=embed_model)
        embedded.append((child, embedding, {"chunk_index": index, "parent_content": parent}))
    return embedded


async def embed_chunks_enriched(
    chunks: list[tuple[str, str]],
    embed_model: str | None = None,
) -> list[tuple[str, list[float], dict]]:
    """
    Generates semantic hypothetical questions and summaries for chunks concurrently 
    using a Semaphore, then embeds the enriched text.
    """
    ollama = OllamaClient()
    semaphore = asyncio.Semaphore(3)  # Limits parallel Ollama calls to protect local CPU
    embedded: list[tuple[str, list[float], dict]] = []

    async def process_single_chunk(index: int, child: str, parent: str):
        async with semaphore:
            # 1. Ask Ollama to generate questions and summary
            enrichment_prompt = (
                "Analyze the text chunk and write 3 short hypothetical questions "
                "that this text directly answers, and a 1-sentence summary. "
                "Return ONLY a JSON object with keys 'questions' (list of strings) and 'summary' (string).\n\n"
                f"Chunk:\n{child}"
            )
            
            questions_list = []
            summary = ""
            try:
                res = await ollama.generate(enrichment_prompt, format="json")
                data = json.loads(res)
                questions_list = data.get("questions") or []
                summary = data.get("summary") or ""
            except Exception as e:
                logger.warning(f"Enrichment generation failed for chunk {index}: {e}")

            questions_str = " ".join(questions_list)
            
            # 2. Build the enriched search content (Intents + Facts)
            enriched_search_text = (
                f"Questions: {questions_str}\n"
                f"Summary: {summary}\n"
                f"Content: {child}"
            )
            
            # 3. Generate embedding vector
            try:
                embedding = await ollama.embed(enriched_search_text, model_name=embed_model)
            except Exception as e:
                logger.error(f"Embedding generation failed for chunk {index}: {e}")
                # Fallback to embedding child raw content if enrichment fails
                embedding = await ollama.embed(child, model_name=embed_model)

            # 4. Save metadata dictionary for UI citations
            metadata = {
                "chunk_index": index,
                "parent_content": parent,
                "hypothetical_questions": questions_list,
                "summary": summary
            }
            return child, embedding, metadata

    # Process chunks concurrently
    tasks = [process_single_chunk(i, child, parent) for i, (child, parent) in enumerate(chunks)]
    results = await asyncio.gather(*tasks)
    
    for res in results:
        if res:
            embedded.append(res)
            
    return embedded
